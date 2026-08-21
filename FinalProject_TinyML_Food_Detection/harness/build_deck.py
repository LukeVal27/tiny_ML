#!/usr/bin/env python3
"""
Build EE446_final_presentation.pptx from scratch.

Rebuilt rather than patched: the old deck carried three superseded live-camera
narratives, and editing them in place left stale numbers in shapes nobody
re-read. Every figure below is transcribed from a file in results/ -- the
provenance comment on each block names which one.

Structure follows the six mandatory topics in the presentation rubric, split
across 10 slides (the rubric explicitly permits exceeding six, and forbids
overcrowding). Slide 6 deliberately carries NO latency/RAM/flash figures: the
rubric bars them from the baseline slide.

    python3 harness/build_deck.py
"""

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = pathlib.Path(__file__).resolve().parent.parent
OUT = REPO / "EE446_final_presentation.pptx"

INK = RGBColor(0x16, 0x18, 0x1D)
MUTE = RGBColor(0x6B, 0x72, 0x80)
FAINT = RGBColor(0xD1, 0xD5, 0xDB)
WASH = RGBColor(0xF6, 0xF7, 0xF9)
ACC = RGBColor(0xC2, 0x41, 0x0C)
GOOD = RGBColor(0x0F, 0x76, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.6)                      # page margin
CW = W - 2 * M                       # content width


# --------------------------------------------------------------------- atoms

def text(slide, x, y, w, h, runs, size=14, color=INK, bold=False,
         align=PP_ALIGN.LEFT, space=6, anchor=MSO_ANCHOR.TOP, line=None):
    """
    `runs` is a string, or a list of paragraphs; a paragraph is a string or a
    list of (text, {overrides}) tuples so one line can mix weights and colours.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if line:
            p.line_spacing = line
        for chunk in (para if isinstance(para, list) else [(para, {})]):
            t, ov = chunk if isinstance(chunk, tuple) else (chunk, {})
            r = p.add_run()
            r.text = t
            r.font.name = FONT
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.color.rgb = ov.get("color", color)
    return tb


def box(slide, x, y, w, h, fill=WASH, edge=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.04
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if edge:
        s.line.color.rgb = edge
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.text_frame.text = ""
    return s


def rule(slide, x, y, w, color=FAINT, weight=1.0):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def slide(prs, title, kicker=None, num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    text(s, M, Inches(0.44), CW, Inches(0.5), title, size=27, bold=True)
    y = Inches(0.95)
    if kicker:
        text(s, M, y, CW, Inches(0.3), kicker, size=13, color=MUTE)
        y = Inches(1.30)
    rule(s, M, y, CW)
    if num:
        text(s, W - M - Inches(0.6), Inches(0.5), Inches(0.6), Inches(0.3),
             num, size=11, color=FAINT, align=PP_ALIGN.RIGHT)
    return s, y + Inches(0.26)


def table(slide, x, y, w, rows, widths=None, fs=12, rh=Inches(0.32),
          highlight=None, align_r=None):
    """Plain table: header rule, no gridlines, optional highlighted body row."""
    nr, nc = len(rows), len(rows[0])
    shp = slide.shapes.add_table(nr, nc, x, y, w, rh * nr)
    tbl = shp.table
    tbl.first_row = False
    tbl.horz_banding = False
    if widths:
        tot = sum(widths)
        for i, fr in enumerate(widths):
            tbl.columns[i].width = Emu(int(w * fr / tot))
    align_r = align_r or set(range(1, nc))

    for ri, row in enumerate(rows):
        tbl.rows[ri].height = rh
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.text = str(val)
            c.margin_left = c.margin_right = Inches(0.06)
            c.margin_top = c.margin_bottom = 0
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            if ri == 0:
                c.fill.fore_color.rgb = WASH
            elif highlight is not None and ri == highlight:
                c.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xEA)
            else:
                c.fill.fore_color.rgb = WHITE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if ci in align_r else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(fs)
                r.font.bold = ri == 0 or (highlight is not None and ri == highlight)
                r.font.color.rgb = (
                    MUTE if ri == 0 else
                    ACC if (highlight is not None and ri == highlight) else INK)
    return shp


def stats(slide, y, items, w=None, gap=Inches(0.25)):
    """A row of number-over-label tiles."""
    w = w or CW
    n = len(items)
    tw = int((w - gap * (n - 1)) / n)
    for i, (val, lab) in enumerate(items):
        x = M + i * (tw + gap)
        text(slide, x, y, Emu(tw), Inches(0.42), val, size=25, bold=True,
             color=ACC)
        text(slide, x, y + Inches(0.44), Emu(tw), Inches(0.3), lab, size=11,
             color=MUTE)


def plain(para):
    """Flatten one paragraph of the runs format back to a bare string."""
    if isinstance(para, str):
        return para
    return "".join(c[0] if isinstance(c, tuple) else c for c in para)


def est_height(body, w_emu, size, space=5, line=1.25):
    """
    Estimate rendered height so a card can size itself to its text.

    Fixed card heights were guesswork and every long body overflowed its box in
    the LibreOffice render. Arial averages ~0.50 em per character at these
    sizes; the 0.93 fudge is measured against the actual renders, not derived.
    """
    total = 0.0
    for para in (body if isinstance(body, list) else [body]):
        # A run may override the paragraph size (a big stat number, say); the
        # tallest run sets the line height, so measure with that.
        sz = size
        if isinstance(para, list):
            sz = max([c[1].get("size", size) for c in para
                      if isinstance(c, tuple)] + [size])
        cpl = max(int((w_emu / 914400) * 96 / (sz * 0.545)), 12)
        lines = max(1, -(-len(plain(para)) // cpl))
        total += lines * sz * line + space
    return Inches(total * 1.10 / 72.0)


def card(slide, x, y, w, min_h, head, body, tint=WASH, hcolor=INK, size=12):
    """
    Draw a card that grows to fit its text; `min_h` is a floor, not a fixed
    height, so side-by-side cards still line up. Returns the bottom edge —
    stack vertically off the return value, never off a hard-coded offset.
    """
    pad = Inches(0.24)
    inner = w - 2 * pad
    # A long header wraps to two lines; the body has to start below wherever it
    # actually ended, not at a fixed offset.
    hd = est_height(head, inner, 13, space=0)
    top = Inches(0.20) + hd + Inches(0.16)
    h = max(min_h, top + est_height(body, inner, size) + pad)
    box(slide, x, y, w, h, fill=tint)
    text(slide, x + pad, y + Inches(0.20), inner, hd, head,
         size=13, bold=True, color=hcolor)
    text(slide, x + pad, y + top, inner, h - top - pad, body,
         size=size, color=INK, space=5, line=1.25)
    return y + h


def note(slide, y, body, tint=WASH, h=Inches(0.72)):
    box(slide, M, y, CW, h, fill=tint)
    text(slide, M + Inches(0.24), y + Inches(0.16), CW - Inches(0.48),
         h - Inches(0.3), body, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE,
         line=1.25)


# -------------------------------------------------------------------- slides

def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # =========================================================== 1 · title
    # Rubric topic 1: title, team, one-sentence objective.
    s = prs.slides.add_slide(prs.slide_layouts[6])
    text(s, M, Inches(1.62), CW, Inches(1.0),
         [[("Food Detection", {}), (" & ", {"color": MUTE}),
           ("Portion Estimation", {})]], size=44, bold=True, space=0)
    text(s, M, Inches(2.70), Inches(11.4), Inches(0.5),
         "Two questions from one camera frame — what food, and how much — "
         "entirely on a $30 microcontroller.", size=17, color=MUTE)
    rule(s, M, Inches(3.42), CW)
    stats(s, Inches(3.75), [
        ("5 + 3", "classes + portion tiers"),      # model/heads.py
        ("216.7 KB", "int8 model on device"),      # results/quantization_real_cw.json
        ("1.08 s", "on-device inference"),         # results/device_runs.jsonl
        ("0.643", "macro-F1, int8 test set"),      # results/quantization_real_cw.json
    ])
    rule(s, M, Inches(5.35), CW)
    text(s, M, Inches(5.62), CW, Inches(0.4),
         [[("Luke Valerio · Daniel Yang", {"bold": True}),
           ("     EE 446 TinyML · Summer 2026 · University of Washington",
            {"color": MUTE})]], size=14)
    text(s, M, Inches(6.15), CW, Inches(0.35),
         "Arduino Nano 33 BLE Sense (nRF52840) · OV7675 camera · TensorFlow "
         "Lite Micro, int8", size=13, color=MUTE)

    # ============================================ 2 · problem + pipeline
    # Rubric topic 2 (part 1): task, sensor, outputs, end-to-end pipeline.
    s, y = slide(prs, "Problem, inputs, and the end-to-end pipeline",
                 "one 96×96 frame in, two predictions out — no network, no host",
                 "01")
    stage = [("OV7675", "176×144\nRGB565"), ("Crop", "144×144\ncentre square"),
             ("Downsample", "96×96×3\n2×2 box mean"), ("Quantise", "int8\nscale 1/255"),
             ("CNN", "157,512\nparameters"), ("Two heads", "5 classes\n3 tiers"),
             ("Output", "Serial · LED\ndashboard")]
    n = len(stage)
    gap = Inches(0.12)
    bw = int((CW - gap * (n - 1)) / n)
    for i, (head, sub) in enumerate(stage):
        x = M + i * (bw + gap)
        tint = RGBColor(0xFD, 0xF2, 0xEA) if i in (4, 5) else WASH
        box(s, x, y, Emu(bw), Inches(1.15), fill=tint)
        text(s, x, y + Inches(0.22), Emu(bw), Inches(0.25), head, size=12,
             bold=True, align=PP_ALIGN.CENTER,
             color=ACC if i in (4, 5) else INK)
        text(s, x, y + Inches(0.55), Emu(bw), Inches(0.5), sub.split("\n"),
             size=10, color=MUTE, align=PP_ALIGN.CENTER, space=0)

    y2 = y + Inches(1.45)
    cw3 = int((CW - Inches(0.3) * 2) / 3)
    card(s, M, y2, Emu(cw3), Inches(1.95), "Task",
         ["Single plate, single food.",
          [("Head A", {"bold": True}), "  5-way classification"],
          [("Head B", {"bold": True}), "  3-tier ordinal portion"],
          "Tiers map to mass: <80 g · 80–180 g · >180 g."])
    card(s, M + Emu(cw3) + Inches(0.3), y2, Emu(cw3), Inches(1.95), "Sensor",
         ["OV7675 on the Tiny ML shield.",
          "QCIF 176×144 RGB565 at 5 fps.",
          "Downsampled, never centre-cropped — the portion tier is "
          "measured against the plate rim, so the whole plate must stay "
          "in frame."])
    card(s, M + 2 * (Emu(cw3) + Inches(0.3)), y2, Emu(cw3), Inches(1.95),
         "Evaluation metrics",
         [[("Class", {"bold": True}), "  macro-F1 (primary), accuracy"],
          [("Portion", {"bold": True}), "  accuracy, mean ordinal error, "
           "off-by-two rate"],
          "Macro-F1 because per-class confusability, not frequency, is the "
          "failure mode."])
    note(s, Inches(6.35),
         [[("Everything runs on the device. ", {"bold": True}),
           ("No companion phone, no cloud inference, no network — the board "
            "captures, preprocesses, infers and reports on its own.", {})]])

    # ======================================================== 3 · the data
    # Rubric topic 2 (part 2). Counts: results/foodseg103_counts.json
    s, y = slide(prs, "Data: real photographed food, honestly labelled",
                 "FoodSeg103 (Apache-2.0) mask cutouts composited onto a "
                 "measured plate", "02")
    table(s, M, y, Inches(7.4), [
        ["our class", "FoodSeg103 source", "cutouts", "caveat"],
        ["broccoli", "broccoli (87)", "986", "clean"],
        ["rice", "rice (66)", "662", "clean"],
        ["potato", "potato (70)", "1,073", "french fries excluded"],
        ["chicken", "chicken duck (48)", "1,226", "merged label"],
        ["beef", "steak (46)", "1,049", "substitution — no beef class"],
    ], widths=[1.1, 1.5, 0.7, 1.9], align_r={2})

    text(s, M, y + Inches(2.20), Inches(7.4), Inches(0.3),
         "Portion labels are constructed, not estimated", size=13, bold=True)
    text(s, M, y + Inches(2.55), Inches(7.4), Inches(1.2),
         ["Food is pasted at a sampled fraction f of plate area, so the label "
          "is exact by definition: mass ≈ f × 513 g (22 cm plate, 1.5 cm deep, "
          "0.9 g/cm³).",
          "Plate apparent diameter is jittered to 62–92% of frame width, so raw "
          "pixel area is an ambiguous cue and the model must measure against "
          "the rim."], size=12, color=INK, line=1.25)

    x2 = M + Inches(7.75)
    w2 = CW - Inches(7.75)
    yb = card(s, x2, y, w2, Inches(1.6), "Corpus",
              [[("4,996", {"bold": True, "color": ACC}), "  RGBA mask cutouts"],
               [("18,000", {"bold": True, "color": ACC}), "  composited plates"],
               "Tiers balanced 5,054 / 5,051 / 4,895.",
               "Split 15,000 train / 3,000 test."])
    card(s, x2, yb + Inches(0.18), w2, Inches(2.0), "Preprocessing",
         ["Resize to 96×96, scale to [0,1].",
          "Augmentation: geometry + sensor degradation only.",
          "No photometric jitter — these are real photographs that already "
          "carry real scene lighting; stacking jitter on top crushed batches "
          "to black and destroyed the colour cue."])

    # ====================================================== 4 · the model
    # Rubric topic 3 (part 1). Params from results/checkpoints/real_cw.keras
    s, y = slide(prs, "Model: one backbone, two heads that pool differently",
                 "depthwise-separable CNN, 157,512 parameters, "
                 "13.6 M MACs per frame", "03")
    text(s, M, y, Inches(7.4), Inches(2.0),
         ["96²×3   →  stem s2  →  48²×16",
          "        →  b1–b2    →  24²×32",
          "        →  b3–b4    →  12²×64",
          "        →  b5–b6    →   6²×128",
          "        →  1×1 widen →  6²×192",
          "                    ├── MAX pool 6×6  →  5 classes    (965 params)",
          "                    └── AVG pool 6×6  →  3 tiers      (579 params)"],
         size=13, space=3)
    text(s, M, y + Inches(2.25), Inches(7.4), Inches(0.6),
         "6 inverted-residual blocks · 3 residual adds · 155,968 parameters "
         "shared between the heads. Deployed ops: CONV_2D ×13, "
         "DEPTHWISE_CONV_2D ×6, ADD ×3, MAX_POOL_2D, AVERAGE_POOL_2D, "
         "FULLY_CONNECTED ×2, SOFTMAX ×2.", size=12, color=MUTE, line=1.25)

    # Design lineage. The efficiency primitives are YOLO-Nano's; the detection
    # formulation is what does not fit in 262 KB, so it is what we dropped.
    yl = y + Inches(3.05)
    text(s, M, yl, Inches(7.4), Inches(0.3),
         "Design lineage — where this sits against the YOLO 'nano' family",
         size=13, bold=True)
    table(s, M, yl + Inches(0.36), Inches(7.4), [
        ["", "params", "compute", "deployed size", "target"],
        ["YOLO Nano (2018)", "—", "4.57 B ops", "~4.0 MB", "Jetson-class"],
        ["YOLOv8n", "3.2 M", "8.7 B FLOPs", "~6 MB", "mobile GPU / CPU"],
        ["ours", "157.5 K", "27.2 M FLOPs", "216.7 KB", "Cortex-M4 @ 64 MHz"],
    ], widths=[1.5, 0.9, 1.1, 1.1, 1.5], highlight=3, fs=11)
    text(s, M, yl + Inches(1.75), Inches(7.4), Inches(0.6),
         [[("Kept: ", {"bold": True}),
           ("depthwise-separable convolutions and residual blocks — the same "
            "efficiency toolkit.  ", {}),
           ("Dropped: ", {"bold": True}),
           ("anchors, box regression, objectness, NMS and the multi-scale "
            "neck. Not for accuracy — for RAM. YOLOv8n's input tensor alone "
            "at 640² is 1.2 MB, 4.7× this chip's entire SRAM.", {})]],
         size=11, color=MUTE, line=1.25)

    yb = card(s, x2, y, w2, Inches(2.2), "Why the heads pool differently",
              ["Food covers a minority of the frame.",
               [("Class head max-pools", {"bold": True}),
                " — presence. A max survives a mostly-empty frame."],
               [("Portion head average-pools", {"bold": True}),
                " — extent. An average measures how much."],
               [("One shared global-average pool stalled the classifier at "
                 "~0.34 validation accuracy.", {"color": ACC})]],
              tint=RGBColor(0xFD, 0xF2, 0xEA), hcolor=ACC)
    card(s, x2, yb + Inches(0.18), w2, Inches(1.3), "Export decisions",
         ["Sized pooling, not Global* — exports better-tested TFLM kernels.",
          "Batch-1 clone on export; a None batch emits dynamic-shape ops.",
          "Outputs bound by width (5 / 3), never by index."])

    # ==================================== 5 · approaches + advanced components
    # Rubric topic 3 (part 2). Numbers: results/compare_on_real_test.json
    s, y = slide(prs, "Approaches evaluated, and our advanced components",
                 "four training configurations, all scored on the same "
                 "3,000-image real test set", "04")
    table(s, M, y, Inches(7.4), [
        ["training configuration", "macro-F1", "class acc", "portion acc"],
        ["synthetic corpus only", "0.5074", "0.5403", "0.8290"],
        ["real corpus, unweighted", "0.5859", "0.6393", "0.9233"],
        ["real + synthetic combined", "0.5330", "0.5957", "0.8770"],
        ["real corpus + class weights", "0.6462", "0.6460", "0.9237"],
    ], widths=[2.2, 1.0, 1.0, 1.0], highlight=4)
    text(s, M, y + Inches(1.95), Inches(7.4), Inches(0.9),
         [[("Mixing corpora is worse than real alone", {"bold": True}),
           (" — replace the synthetic data, do not augment with it.", {})],
          [("Class weights are required", {"bold": True}),
           (" — without them chicken collapses to F1 0.115 / recall 0.063 on a "
            "perfectly balanced corpus. That is confusability, not frequency.",
            {})]], size=12, line=1.25)

    yb = card(s, x2, y, w2, Inches(1.7),
              "Advanced component 1 — multiple approaches compared",
              ["Four corpora/weighting configurations, plus a head-pooling "
               "ablation, each trained and scored end to end on an identical "
               "test set.",
               "The corpus choice moved macro-F1 by 0.139 — more than any "
               "architecture change we tried."],
              tint=RGBColor(0xEC, 0xFD, 0xF5), hcolor=GOOD)
    card(s, x2, yb + Inches(0.18), w2, Inches(1.7),
         "Advanced component 2 — output beyond the Serial Monitor",
         ["On-board RGB LED signals capture, result and camera fault without "
          "a host attached.",
          "A live host dashboard (harness/camera_view.py) streams the frame "
          "the model actually sees, with a plate-framing guide ring and the "
          "prediction overlaid."], tint=RGBColor(0xEC, 0xFD, 0xF5), hcolor=GOOD)

    # ============================================ 6 · baseline performance
    # Rubric topic 4. NO latency/RAM/flash on this slide -- rubric forbids it.
    # Numbers: results/quantization_real_cw.json
    s, y = slide(prs, "Baseline model performance — before compression",
                 "float32, 3,000-image held-out test set · no compression "
                 "applied yet", "05")
    table(s, M, y, Inches(7.4), [
        ["split", "macro-F1", "class acc", "portion acc", "ordinal err"],
        ["train", "0.6832", "0.6850", "0.9261", "—"],
        ["validation", "0.6475", "0.6376", "0.9253", "—"],
        ["test", "0.6462", "0.6460", "0.9237", "0.0763"],
    ], widths=[1.3, 1.0, 1.0, 1.0, 1.0], highlight=3)
    text(s, M, y + Inches(1.55), Inches(7.4), Inches(0.5),
         "Train-to-test gap 0.037 — the model is not overfitting. The ceiling "
         "is set by the source data, not by capacity.", size=12, color=MUTE,
         line=1.25)

    text(s, M, y + Inches(2.25), Inches(7.4), Inches(0.3),
         "Per-class F1, float32 test set", size=13, bold=True)
    table(s, M, y + Inches(2.62), Inches(7.4), [
        ["broccoli", "beef", "rice", "potato", "chicken"],
        ["0.969", "0.665", "0.567", "0.522", "0.508"],
    ], widths=[1, 1, 1, 1, 1], align_r=set(range(5)))

    yb = card(s, x2, y, w2, Inches(1.8), "Portion head",
              [[("0.9237", {"bold": True, "size": 20, "color": ACC}),
                ("  tier accuracy", {"color": MUTE})],
               [("0 of 3,000", {"bold": True}), " off-by-two errors — when it "
                "is wrong, it is wrong by exactly one tier."],
               "Mean ordinal error 0.0763."])
    card(s, x2, yb + Inches(0.18), w2, Inches(1.7), "Where the class errors are",
         ["Rice ↔ potato is the dominant confusion: 183 rice→potato and "
          "126 potato→rice.",
          "Both are pale, starchy and near-textureless at 96×96. Broccoli, "
          "which is unmistakably green, reaches 0.969."])

    # ================================= 7 · compressed model + resource cost
    # Rubric topic 5 (part 1). Sizes: results/quantization_real_cw.json
    # Device figures: results/device_runs.jsonl, results/compile_sweep.jsonl
    s, y = slide(prs, "Compression: int8 post-training quantization",
                 "2.8× smaller for 0.003 macro-F1 — the whole TinyML trade in "
                 "one row", "06")
    table(s, M, y, CW, [
        ["model", "macro-F1", "portion acc", "model size", "tensor arena",
         "flash", "latency"],
        ["float32 (baseline)", "0.6462", "0.9237", "621,956 B", "—", "—",
         "does not fit"],
        ["int8 PTQ (deployed)", "0.6432", "0.9187", "221,888 B", "113,516 B",
         "500,648 B", "1,082 ms"],
    ], widths=[1.6, 1.0, 1.1, 1.1, 1.0, 1.0, 1.0], highlight=2, fs=12)

    y3 = y + Inches(1.35)
    stats(s, y3, [
        ("2.80×", "compression ratio"),
        ("−0.0030", "macro-F1 cost"),
        ("43.3%", "of SRAM, arena alone"),
        ("50.9%", "of 983 KB flash"),
        ("1,082 ms", "mean of 20 runs"),
    ])

    y4 = y3 + Inches(1.15)
    cw2 = int((CW - Inches(0.3)) / 2)
    card(s, M, y4, Emu(cw2), Inches(1.75), "Quantization settings",
         ["Full-integer post-training quantization; int8 weights and "
          "activations, int8 input and output.",
          "Representative dataset: 300 training images with augmentation ON — "
          "the quantiser should see the degraded activations the device will "
          "actually produce."])
    card(s, M + Emu(cw2) + Inches(0.3), y4, Emu(cw2), Inches(1.75),
         "Measured RAM budget, on hardware",
         ["Fixed overhead (mbed OS + TFLM + Serial)  51,792 B",
          "RGB565 QCIF frame buffer  50,688 B",
          "Tensor arena in use  113,516 B",
          "Total static RAM  222,688 B of 262,144 B",
          "Largest free block still allocatable  54,284 B"])

    # ======================================== 8 · on-device evaluation
    # Rubric topic 5 (part 2): sample counts MUST be explicit.
    # Source: results/live_captures.jsonl, session=final
    s, y = slide(prs, "On-device evaluation — live camera, labelled captures",
                 "real plates photographed by the OV7675 · sample counts "
                 "stated explicitly", "07")
    table(s, M, y, Inches(7.4), [
        ["class", "# on-device samples", "correct", "recall", "F1"],
        ["beef", "23", "17", "0.739", "0.850"],
        ["chicken", "20", "6", "0.300", "0.387"],
        ["broccoli", "19", "15", "0.789", "0.833"],
        ["rice", "13", "12", "0.923", "0.600"],
        ["potato", "0", "—", "food not available", ""],
        ["overall", "75", "50", "0.667", "0.668"],
    ], widths=[1.1, 1.5, 0.8, 1.2, 0.7], highlight=6)

    text(s, M, y + Inches(2.55), Inches(7.4), Inches(0.6),
         "75 labelled captures with harness/live_capture.py. Macro-F1 is over "
         "the four classes with on-device support; potato was not available "
         "and is reported as an explicit zero.", size=12, color=MUTE,
         line=1.25)

    # Predicted before the session, then confirmed by it -- stated in that
    # order, because it is a prediction and not a rationalisation.
    text(s, M, y + Inches(3.20), Inches(7.4), Inches(1.1),
         [[("Our chicken and beef are raw; the training food is cooked. "
            "We predicted the effect before capturing, and it cut both ways.",
            {"bold": True})],
          [("Pale raw chicken collides with rice", {"bold": True}),
           (" — 11 of 20 chicken captures called rice, 3 potato.  ", {}),
           ("Raw beef goes the other way", {"bold": True}),
           (" — the measured failure was that the camera renders cooked beef "
            "neutral grey (R−B +1.4 against +39.8 in training). Raw beef is "
            "strongly red, and recall went from 0/5 to 17/23 at precision "
            "1.000.", {})]],
         size=12, line=1.25)

    yb = card(s, x2, y, w2, Inches(1.5), "Capture-to-result latency, measured",
              ["Capture (bit-banged readFrame)  331–583 ms",
               "Preprocess (crop, downsample, quantise)  ~68 ms",
               "Inference  1,082 ms",
               [("Total  1,482–1,733 ms", {"bold": True, "color": ACC})]])
    yb = card(s, x2, yb + Inches(0.18), w2, Inches(1.5),
              "Portion head on device — a degenerate output",
              [[("The head emitted 'large' on all 75 captures", {"bold": True}),
                (", every one at exactly 0.9961 = 255/256, the int8 softmax "
                 "ceiling. It never once produced 'small' or 'medium'.", {})],
               [("Apparent accuracy 0.478 (33/69 labelled) is the base rate "
                 "alone — 33 of the plates happened to be large. Recall 1.000 "
                 "on large, 0.000 on small and medium.", {"color": ACC})]],
              tint=RGBColor(0xFD, 0xF2, 0xEA), hcolor=ACC)
    # The Invoke() aliasing bug is told in full on the findings slide; here it
    # only needs the one line that says the class head is trustworthy.
    text(s, x2, yb + Inches(0.26), w2, Inches(0.5),
         [[("Class head verified against the host: ", {"bold": True}),
           ("img_sum = input_sum = 449,425 · match = OK", {})]],
         size=12, color=MUTE, line=1.25)

    # ============================================== 9 · key findings
    # Rubric topic 6 (part 1).
    s, y = slide(prs, "Key findings", "what moved the numbers, and what did not",
                 "08")
    cw3b = int((CW - Inches(0.3) * 2) / 3)
    card(s, M, y, Emu(cw3b), Inches(2.25),
         "Data beat architecture, decisively",
         ["Switching from synthetic textures to real FoodSeg103 cutouts, plus "
          "class weighting, moved macro-F1 from 0.5074 to 0.6462.",
          [("+0.139 from data alone", {"bold": True, "color": ACC}),
           (", at identical model size, arena and latency.", {})]])
    card(s, M + Emu(cw3b) + Inches(0.3), y, Emu(cw3b), Inches(2.25),
         "int8 was nearly free",
         ["2.80× smaller for 0.0030 macro-F1 and 0.005 portion accuracy.",
          "The portion head lost nothing structurally: still zero off-by-two "
          "errors in 3,000 test images after quantization."])
    card(s, M + 2 * (Emu(cw3b) + Inches(0.3)), y, Emu(cw3b), Inches(2.25),
         "The bottleneck is not the kernels",
         ["CMSIS-NN is already active — 26 int8 symbols, 251 DSP SIMD "
          "instructions — and measured no speedup at all "
          "(1,092 ms vs 1,076 ms reference).",
          "1,082 ms remains unexplained."])

    y5 = y + Inches(2.45)
    card(s, M, y5, Emu(cw2), Inches(1.85),
         "Challenge — the device disagreed with the host",
         ["A byte-identical input produced a different class on device. Three "
          "hypotheses were wrong before an on-device checksum isolated it: "
          "Invoke() aliases and destroys its own input tensor.",
          [("Fix: refill the input before every Invoke().", {"bold": True})]],
         tint=RGBColor(0xFD, 0xF2, 0xEA), hcolor=ACC)
    card(s, M + Emu(cw2) + Inches(0.3), y5, Emu(cw2), Inches(1.85),
         "Challenge — the portion head pins at 'large' on camera input",
         ["Diagnosed on the host, not guessed. Ruled out: white balance, "
          "contrast, and mean/std matching all leave it pinned.",
          [("A scale sweep shows the head is correct from 0.45× to 1.35× "
            "plate zoom, and only fails once the rim leaves the frame.",
            {"bold": True})]],
         tint=RGBColor(0xFD, 0xF2, 0xEA), hcolor=ACC)

    # ======================================= 10 · limitations + future
    # Rubric topic 6 (part 2).
    s, y = slide(prs, "Limitations and what we would do next",
                 "the honest boundary of what these numbers support", "09")
    yl1 = card(s, M, y, Emu(cw2), Inches(2.55),
               "Immovable without new data",
         [[("FoodSeg103 has no beef and no pure chicken.", {"bold": True}),
           (" Our beef is steak; our chicken is a merged 'chicken duck' label. "
            "This caps chicken near F1 0.51 — training cannot fix it.", {})],
          [("Rice and potato are not separable at 96×96.", {"bold": True}),
           (" Grain structure sits at the resolution limit, and higher "
            "resolution breaks the RAM budget.", {})],
          [("The food is real but the plate is drawn.", {"bold": True}),
           (" Real shadows, specular highlights and depth are absent from "
            "training.", {})]])
    yl2 = card(s, M + Emu(cw2) + Inches(0.3), y, Emu(cw2), Inches(2.55),
               "Open and unverified",
         [[("Portion tiers are modelled, not weighed.", {"bold": True}),
           (" 0.9187 is scored against our own definition of a tier; "
            "real gram accuracy is unmeasured.", {})],
          [("The portion head is degenerate on real camera input.", {"bold": True}),
           (" It emitted 'large' on all 75 captures and never produced another "
            "tier. Host diagnosis rules out white balance, contrast and "
            "mean/std matching; a scale sweep shows the head is correct from "
            "0.45× to 1.35× plate zoom. Unresolved.", {})],
          [("Single-label only.", {"bold": True}),
           (" The class head ends in softmax, so it reports one food per "
            "plate.", {})]])

    note(s, max(yl1, yl2) + Inches(0.20),
         [[("Next: ", {"bold": True}),
           ("multi-label is feasible on this hardware at essentially zero "
            "memory cost — 5 independent sigmoids plus a 5×3 per-class tier "
            "head takes the output tensor from 8 to 20 values and leaves the "
            "arena unchanged. The class head already max-pools, the correct "
            "operator for presence detection, so only the activation, loss and "
            "data generator change. Beyond that: ~150 weighed captures at a "
            "fixed camera height would close the plate-realism and gram-"
            "accuracy gaps together.", {})]], h=Inches(1.15))

    text(s, M, Inches(6.95), CW, Inches(0.3),
         "Luke Valerio · Daniel Yang     |     EE 446 TinyML, Summer 2026",
         size=11, color=MUTE)

    prs.save(OUT)
    print(f"wrote {OUT.name}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
